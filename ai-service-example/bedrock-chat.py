import streamlit as st # type: ignore
import boto3
from botocore.config import Config
import os
import pandas as pd
import time
import json
import io
import re
import openpyxl
from python_calamine import CalamineWorkbook
from openpyxl.cell import Cell
from openpyxl.worksheet.cell_range import CellRange
from docx.table import _Cell
from pptx import Presentation
from botocore.exceptions import ClientError
import pytesseract
from PIL import Image
import PyPDF2
import chardet
from docx import Document as DocxDocument
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from docx.document import Document
from docx.text.paragraph import Paragraph
from docx.table import Table as DocxTable
import concurrent.futures
from functools import partial
import textract
import random
import psycopg2
from psycopg2.extras import Json

config = Config(
    read_timeout=600,  # Read timeout parameter
    retries=dict(
        max_attempts=10  # Handle retries
    )
)
import function_calling_utils
st.set_page_config(initial_sidebar_state="auto")

# Read credentials
with open('config.json', 'r', encoding='utf-8') as f:
    config_file = json.load(f)
# pricing info
with open('pricing.json', 'r', encoding='utf-8') as f:
    pricing_file = json.load(f)

S3 = boto3.client('s3')
COGNITO = boto3.client('cognito-idp')
LOCAL_CHAT_FILE_NAME = "chat-history.json"
BUCKET = config_file["Bucket_Name"]
OUTPUT_TOKEN = config_file.get("max-output-token", 2000) 
S3_DOC_CACHE_PATH = config_file["document-upload-cache-s3-path"]
LOAD_DOC_IN_ALL_CHAT_CONVO = config_file["load-doc-in-chat-history"]
CHAT_HISTORY_LENGTH = config_file["chat-history-loaded-length"]
REGION = config_file["bedrock-region"]
CSV_SEPERATOR = config_file["csv-delimiter"]
INPUT_BUCKET = config_file["input_bucket"]
INPUT_S3_PATH = config_file["input_s3_path"]
INPUT_EXT = tuple(f".{x}" for x in config_file["input_file_ext"].split(','))

# PostgreSQL configuration
conn = psycopg2.connect(
    host=config_file["postgres_host"],
    database=config_file["postgres_database"],
    user=config_file["postgres_user"],
    password=config_file["postgres_password"]
)
cursor = conn.cursor()

bedrock_runtime = boto3.client(service_name='bedrock-runtime', region_name=REGION, config=config)

if 'messages' not in st.session_state:
    st.session_state['messages'] = []
if 'input_token' not in st.session_state:
    st.session_state['input_token'] = 0
if 'output_token' not in st.session_state:
    st.session_state['output_token'] = 0
if 'chat_hist' not in st.session_state:
    st.session_state['chat_hist'] = []
if 'user_sess' not in st.session_state:
    st.session_state['user_sess'] = str(time.time())
if 'chat_session_list' not in st.session_state:
    st.session_state['chat_session_list'] = []
if 'count' not in st.session_state:
    st.session_state['count'] = 0
if 'userid' not in st.session_state:
    st.session_state['userid'] = config_file["UserId"]
if 'cost' not in st.session_state:
    st.session_state['cost'] = 0

def get_object_with_retry(bucket, key):
    max_retries = 5
    retries = 0
    backoff_base = 2
    max_backoff = 3  # Maximum backoff time in seconds
    s3 = boto3.client('s3')
    while retries < max_retries:
        try:
            response = s3.get_object(Bucket=bucket, Key=key)
            return response
        except ClientError as e:
            error_code = e.response['Error']['Code']
            if error_code == 'DecryptionFailureException':
                sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                print(f"Decryption failed, retrying in {sleep_time} seconds...")
                time.sleep(sleep_time)
                retries += 1
            elif e.response['Error']['Code'] == 'ModelStreamErrorException':
                if retries < max_retries:
                    # Throttling, exponential backoff
                    sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                    time.sleep(sleep_time)
                    retries += 1
            else:
                raise e

    # If we reach this point, it means the maximum number of retries has been exceeded
    raise Exception(f"Failed to get object {key} from bucket {bucket} after {max_retries} retries.")

def save_chat_local(file_path, new_data, session_id):
    """Store long term chat history on local disk"""
    try:
        # Read the existing JSON data from the file
        with open(file_path, "r", encoding='utf-8') as file:
            existing_data = json.load(file)
        if session_id not in existing_data:
            existing_data[session_id] = []
    except FileNotFoundError:
        # If the file doesn't exist, initialize an empty list
        existing_data = {session_id: []}
    # Append the new data to the existing list
    from decimal import Decimal
    data = [{k: float(v) if isinstance(v, Decimal) else v for k, v in item.items()} for item in new_data]
    existing_data[session_id].extend(data)

    # Write the updated list back to the JSON file
    with open(file_path, "w", encoding="utf-8") as file:
        json.dump(existing_data, file)

def load_chat_local(file_path, session_id):
    """Load long term chat history from local disk"""
    try:
        # Read the existing JSON data from the file
        with open(file_path, "r", encoding='utf-8') as file:
            existing_data = json.load(file)
            if session_id in existing_data:
                existing_data = existing_data[session_id]
            else:
                existing_data = []
    except FileNotFoundError:
        # If the file doesn't exist, initialize an empty list
        existing_data = []
    return existing_data

def process_files(files):
    result_string = ""
    errors = []
    future_proxy_mapping = {}
    futures = []

    with concurrent.futures.ProcessPoolExecutor() as executor:
        # Partial function to pass the handle_doc_upload_or_s3 function
        func = partial(handle_doc_upload_or_s3)
        for file in files:
            future = executor.submit(func, file)
            future_proxy_mapping[future] = file
            futures.append(future)

        # Collect the results and handle exceptions
        for future in concurrent.futures.as_completed(futures):
            file_url = future_proxy_mapping[future]
            try:
                result = future.result()
                doc_name = os.path.basename(file_url)

                result_string += f"<{doc_name}>\n{result}\n</{doc_name}>\n"
            except Exception as e:
                # Get the original function arguments from the Future object
                error = {'file': file_url, 'error': str(e)}
                errors.append(error)

    return errors, result_string

def handle_doc_upload_or_s3(file, cutoff=None):
    """Handle various document formats"""
    dir_name, ext = os.path.splitext(file)
    if ext.lower() in [".pdf", ".png", ".jpg", ".tif", ".jpeg"]:
        content = exract_pdf_text(file)
    elif ".csv" == ext.lower():
        content = parse_csv_from_s3(file, cutoff)
    elif ext.lower() in [".xlsx", ".xls"]:
        content = table_parser_utills(file, cutoff)
    elif ".json" == ext.lower():
        obj = get_s3_obj_from_bucket_(file)
        content = json.loads(obj['Body'].read())
    elif ext.lower() in [".txt", ".py"]:
        obj = get_s3_obj_from_bucket_(file)
        content = obj['Body'].read()
    elif ".docx" == ext.lower():
        obj = get_s3_obj_from_bucket_(file)
        content = obj['Body'].read()
        docx_buffer = io.BytesIO(content)
        content = extract_text_and_tables(docx_buffer)
    elif ".pptx" == ext.lower():
        obj = get_s3_obj_from_bucket_(file)
        content = obj['Body'].read()
        docx_buffer = io.BytesIO(content)
        content = extract_text_from_pptx_s3(docx_buffer)
    else:
        obj = get_s3_obj_from_bucket_(file)
        content = obj['Body'].read()
        doc_buffer = io.BytesIO(content)
        content = textract.process(doc_buffer).decode()
    # Implement any other file extension logic
    return content

class InvalidContentError(Exception):
    pass

def detect_encoding(s3_uri):
    """Detect CSV encoding"""
    s3 = boto3.client('s3')
    match = re.match("s3://(.+?)/(.+)", s3_uri)
    if match:
        bucket_name = match.group(1)
        key = match.group(2)
    response = s3.get_object(Bucket=bucket_name, Key=key)
    content = response['Body'].read()
    result = chardet.detect(content)
    return result['encoding']

def parse_csv_from_s3(s3_uri, cutoff):
    """Read CSV files"""
    try:
        # Detect the file encoding using chardet
        encoding = detect_encoding(s3_uri)
        # Sniff the delimiter and read the CSV file
        df = pd.read_csv(s3_uri, delimiter=None, engine='python', encoding=encoding)
        if cutoff:
            df = df.iloc[:20]
        return df.to_csv(index=False, sep=CSV_SEPERATOR)
    except Exception as e:
        raise InvalidContentError(f"Error: {e}")

def iter_block_items(parent):
    if isinstance(parent, Document):
        parent_elm = parent.element.body
    elif isinstance(parent, _Cell):
        parent_elm = parent._tc
    else:
        raise ValueError("something's not right")

    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield DocxTable(child, parent)

def extract_text_and_tables(docx_path):
    """Extract text from docx files"""
    document = DocxDocument(docx_path)
    content = ""
    current_section = ""
    section_type = None
    for block in iter_block_items(document):
        if isinstance(block, Paragraph):
            if block.text:
                if block.style.name == 'Heading 1':
                    # Close the current section if it exists
                    if current_section:
                        content += f"{current_section}</{section_type}>\n"
                        current_section = ""
                        section_type = None
                    section_type = "h1"
                    content += f"<{section_type}>{block.text}</{section_type}>\n"
                elif block.style.name == 'Heading 3':
                    # Close the current section if it exists
                    if current_section:
                        content += f"{current_section}</{section_type}>\n"
                        current_section = ""
                    section_type = "h3"
                    content += f"<{section_type}>{block.text}</{section_type}>\n"

                elif block.style.name == 'List Paragraph':
                    # Add to the current list section
                    if section_type != "list":
                        # Close the current section if it exists
                        if current_section:
                            content += f"{current_section}</{section_type}>\n"
                        section_type = "list"
                        current_section = "<list>"
                    current_section += f"{block.text}\n"
                elif block.style.name.startswith('toc'):
                    # Add to the current toc section
                    if section_type != "toc":
                        # Close the current section if it exists
                        if current_section:
                            content += f"{current_section}</{section_type}>\n"
                        section_type = "toc"
                        current_section = "<toc>"
                    current_section += f"{block.text}\n"
                else:
                    # Close the current section if it exists
                    if current_section:
                        content += f"{current_section}</{section_type}>\n"
                        current_section = ""
                        section_type = None

                    # Append the passage text without tagging
                    content += f"{block.text}\n"

        elif isinstance(block, DocxTable):
            # Add the current section before the table
            if current_section:
                content += f"{current_section}</{section_type}>\n"
                current_section = ""
                section_type = None

            content += "<table>\n"
            for row in block.rows:
                row_content = []
                for cell in row.cells:
                    cell_content = []
                    for nested_block in iter_block_items(cell):
                        if isinstance(nested_block, Paragraph):
                            cell_content.append(nested_block.text)
                        elif isinstance(nested_block, DocxTable):
                            nested_table_content = parse_nested_table(nested_block)
                            cell_content.append(nested_table_content)
                    row_content.append(CSV_SEPERATOR.join(cell_content))
                content += CSV_SEPERATOR.join(row_content) + "\n"
            content += "</table>\n"

    # Add the final section
    if current_section:
        content += f"{current_section}</{section_type}>\n"

    return content

def parse_nested_table(table):
    nested_table_content = "<table>\n"
    for row in table.rows:
        row_content = []
        for cell in row.cells:
            cell_content = []
            for nested_block in iter_block_items(cell):
                if isinstance(nested_block, Paragraph):
                    cell_content.append(nested_block.text)
                elif isinstance(nested_block, DocxTable):
                    nested_table_content += parse_nested_table(nested_block)
            row_content.append(CSV_SEPERATOR.join(cell_content))
        nested_table_content += CSV_SEPERATOR.join(row_content) + "\n"
    nested_table_content += "</table>"
    return nested_table_content

def extract_text_from_pptx_s3(pptx_buffer):
    """Extract text from pptx files"""
    presentation = Presentation(pptx_buffer)
    text_content = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                slide_text.append(shape.text)
        text_content.append('\n'.join(slide_text))
    return '\n\n'.join(text_content)

def exract_pdf_text(file):
    file_base_name = os.path.basename(file)
    dir_name, ext = os.path.splitext(file)
    s3 = boto3.resource("s3")
    match = re.match("s3://(.+?)/(.+)", file)
    if match:
        bucket_name = match.group(1)
        key = match.group(2)
    if "pdf" in ext.lower():
        pdf_bytes = io.BytesIO()
        s3.Bucket(bucket_name).download_fileobj(key, pdf_bytes)
        pdf_bytes.seek(0)
        pdf_reader = PyPDF2.PdfReader(pdf_bytes)
        text = ''
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            page_text = page.extract_text()
            if page_text:
                text += page_text
            else:
                # If the page has no text (could be an image), use OCR
                page_image = page_to_image(page)
                text += pytesseract.image_to_string(page_image)
    else:
        img_bytes = io.BytesIO()
        s3.Bucket(bucket_name).download_fileobj(key, img_bytes)
        img_bytes.seek(0)
        image = Image.open(img_bytes)
        text = pytesseract.image_to_string(image)
    return text

def page_to_image(page):
    """Convert a PDF page to an image for OCR processing"""
    # You might need to use a library like pdf2image
    from pdf2image import convert_from_bytes
    images = convert_from_bytes(page.extract_text())
    return images[0]

def strip_newline(cell):
    return str(cell).strip()

def table_parser_openpyxl(file, cutoff):
    # Read from S3
    s3 = boto3.client('s3')
    match = re.match("s3://(.+?)/(.+)", file)
    if match:
        bucket_name = match.group(1)
        key = match.group(2)
        obj = s3.get_object(Bucket=bucket_name, Key=key)
        # Read Excel file from S3 into a buffer
        xlsx_buffer = io.BytesIO(obj['Body'].read())
        xlsx_buffer.seek(0)
        # Load workbook
        wb = openpyxl.load_workbook(xlsx_buffer)
        all_sheets_string = ""
        # Iterate over each sheet in the workbook
        for sheet_name in wb.sheetnames:
            worksheet = wb[sheet_name]

            all_merged_cell_ranges = list(
                worksheet.merged_cells.ranges
            )
            for merged_cell_range in all_merged_cell_ranges:
                merged_cell = merged_cell_range.start_cell
                worksheet.unmerge_cells(range_string=merged_cell_range.coord)
                for row_index, col_index in merged_cell_range.cells:
                    cell = worksheet.cell(row=row_index, column=col_index)
                    cell.value = merged_cell.value
            # Convert sheet data to a DataFrame
            df = pd.DataFrame(worksheet.values)
            df = df.applymap(strip_newline)
            if cutoff:
                df = df.iloc[:20]

            # Convert to string and tag by sheet name
            tabb = df.to_csv(sep=CSV_SEPERATOR, index=False, header=0)
            all_sheets_string += f'<{sheet_name}>\n{tabb}\n</{sheet_name}>\n'
        return all_sheets_string
    else:
        raise Exception(f"{file} not formatted as an S3 path")

def calamaine_excel_engine(file, cutoff):
    # Read from S3
    s3 = boto3.client('s3')
    match = re.match("s3://(.+?)/(.+)", file)
    if match:
        bucket_name = match.group(1)
        key = match.group(2)
        obj = s3.get_object(Bucket=bucket_name, Key=key)
        # Read Excel file from S3 into a buffer
        xlsx_buffer = io.BytesIO(obj['Body'].read())
        xlsx_buffer.seek(0)
        all_sheets_string = ""
        # Load the Excel file
        workbook = CalamineWorkbook.from_filelike(xlsx_buffer)
        # Iterate over each sheet in the workbook
        for sheet_name in workbook.sheet_names:
            # Get the sheet by name
            sheet = workbook.get_sheet_by_name(sheet_name)
            df = pd.DataFrame(sheet.to_python(skip_empty_area=False))
            df = df.applymap(strip_newline)
            if cutoff:
                df = df.iloc[:20]
            tabb = df.to_csv(sep=CSV_SEPERATOR, index=False, header=0)
            all_sheets_string += f'<{sheet_name}>\n{tabb}\n</{sheet_name}>\n'
        return all_sheets_string
    else:
        raise Exception(f"{file} not formatted as an S3 path")

def table_parser_utills(file, cutoff):
    try:
        response = table_parser_openpyxl(file, cutoff)
        if response:
            return response
        else:
            return calamaine_excel_engine(file, cutoff)
    except Exception as e:
        try:
            return calamaine_excel_engine(file, cutoff)
        except Exception as e:
            raise Exception(str(e))

def put_db(params, messages):
    """Store long term chat history in PostgreSQL"""
    chat_item = {
        "UserID": st.session_state['userid'],  
        "SessionID": params["session_id"],  
        "Messages": [messages],  
    }

    # Check if the record already exists
    cursor.execute("""
        SELECT "Messages" FROM "RangDongChatbot" WHERE "UserID"=%s AND "SessionID"=%s
    """, (st.session_state['userid'], params["session_id"]))
    result = cursor.fetchone()

    if result:
        existing_messages = result[0]
        chat_item["Messages"] = existing_messages + [messages]
        cursor.execute("""
            UPDATE "RangDongChatbot" SET "Messages"=%s WHERE "UserID"=%s AND "SessionID"=%s
        """, (Json(chat_item["Messages"]), st.session_state['userid'], params["session_id"]))
    else:
        cursor.execute("""
            INSERT INTO "RangDongChatbot" ("UserID", "SessionID", "Messages")
            VALUES (%s, %s, %s)
        """, (st.session_state['userid'], params["session_id"], Json(chat_item["Messages"])))

    conn.commit()

def get_chat_history_db(params, cutoff, claude3):
    current_chat, chat_hist = [], []
    cursor.execute("""
        SELECT "Messages" FROM "RangDongChatbot" WHERE "UserID"=%s AND "SessionID"=%s
    """, (st.session_state['userid'], params["session_id"]))
    result = cursor.fetchone()

    if result:
        chat_hist = result[0][-cutoff:]
        for d in chat_hist:
            if d.get('image') and claude3 and LOAD_DOC_IN_ALL_CHAT_CONVO:
                content = []
                for img in d['image']:
                    s3 = boto3.client('s3')
                    match = re.match("s3://(.+?)/(.+)", img)
                    image_name = os.path.basename(img)
                    _, ext = os.path.splitext(image_name)
                    if "jpg" in ext:
                        ext = ".jpeg"
                    bucket_name = match.group(1)
                    key = match.group(2)
                    obj = s3.get_object(Bucket=bucket_name, Key=key)
                    bytes_image = obj['Body'].read()
                    content.extend([{"text": image_name}, {
                        "image": {
                            "format": f"{ext.lower().replace('.', '')}",
                            "source": {"bytes": bytes_image}
                        }
                    }])
                content.extend([{"text": d['user']}])
                current_chat.append({'role': 'user', 'content': content})
            elif d.get('document') and LOAD_DOC_IN_ALL_CHAT_CONVO:
                if 'tool_use_id' in d and d['tool_use_id']:
                    doc = 'Here are the documents:\n'
                    for docs in d['document']:
                        uploads = handle_doc_upload_or_s3(docs, 20)
                        doc_name = os.path.basename(docs)
                        doc += f"<{doc_name}>\n{uploads}\n</{doc_name}>\n"
                else:
                    doc = 'Here are the documents:\n'
                    for docs in d['document']:
                        uploads = handle_doc_upload_or_s3(docs)
                        doc_name = os.path.basename(docs)
                        doc += f"<{doc_name}>\n{uploads}\n</{doc_name}>\n"
                if not claude3 and d.get("image"):
                    for docs in d['image']:
                        uploads = handle_doc_upload_or_s3(docs)
                        doc_name = os.path.basename(docs)
                        doc += f"<{doc_name}>\n{uploads}\n</{doc_name}>\n"
                current_chat.append({'role': 'user', 'content': [{"text": doc + d['user']}]})
            else:
                current_chat.append({'role': 'user', 'content': [{"text": d['user']}]})
            current_chat.append({'role': 'assistant', 'content': [{"text": d['assistant']}]})
    else:
        chat_hist = []
    return current_chat, chat_hist

def get_s3_keys(prefix):
    """List all keys in an S3 path"""
    s3 = boto3.client('s3')
    keys = []
    next_token = None
    while True:
        if next_token:
            response = s3.list_objects_v2(Bucket=BUCKET, Prefix=prefix, ContinuationToken=next_token)
        else:
            response = s3.list_objects_v2(Bucket=BUCKET, Prefix=prefix)
        if "Contents" in response:
            for obj in response['Contents']:
                key = obj['Key']
                name = key[len(prefix):]
                keys.append(name)
        if "NextContinuationToken" in response:
            next_token = response["NextContinuationToken"]
        else:
            break
    return keys

def parse_s3_uri(uri):
    """
    Parse an S3 URI and extract the bucket name and key.
    """
    pattern = r'^s3://([^/]+)/(.*)$'
    match = re.match(pattern, uri)
    if match:
        return match.groups()
    return (None, None)

def copy_s3_object(source_uri, dest_bucket, dest_key):
    """
    Copy an object from one S3 location to another.
    """
    s3 = boto3.client('s3')

    # Parse the source URI
    source_bucket, source_key = parse_s3_uri(source_uri)
    if not source_bucket or not source_key:
        print(f"Invalid source URI: {source_uri}")
        return False

    try:
        # Create a copy source dictionary
        copy_source = {
            'Bucket': source_bucket,
            'Key': source_key
        }
        source_key = os.path.basename(source_key)
        # Copy the object
        s3.copy_object(CopySource=copy_source, Bucket=dest_bucket, Key=f"{dest_key}/{source_key}")

        print(f"File copied from {source_uri} to s3://{dest_bucket}/{dest_key}/{source_key}")
        return f"s3://{dest_bucket}/{dest_key}/{source_key}"

    except ClientError as e:
        print(f"An error occurred: {e}")
        raise e

def get_s3_obj_from_bucket_(file):
    s3 = boto3.client('s3')
    match = re.match("s3://(.+?)/(.+)", file)
    if match:
        bucket_name = match.group(1)
        key = match.group(2)
        obj = s3.get_object(Bucket=bucket_name, Key=key)
    return obj

def put_obj_in_s3_bucket_(docs):
    if isinstance(docs, str):
        s3_uri_pattern = r'^s3://([^/]+)/(.*?([^/]+)/?)$'
        if bool(re.match(s3_uri_pattern, docs)):
            file_uri = copy_s3_object(docs, BUCKET, S3_DOC_CACHE_PATH)
            return file_uri
    else:
        file_name = os.path.basename(docs.name)
        file_path = f"{S3_DOC_CACHE_PATH}/{file_name}"
        S3.put_object(Body=docs.read(), Bucket=BUCKET, Key=file_path)
        return f"s3://{BUCKET}/{file_path}"

def bedrock_streemer(params, response, handler):
    text = ''
    for chunk in response['stream']:

        if 'contentBlockDelta' in chunk:
            delta = chunk['contentBlockDelta']['delta']
            if 'text' in delta:
                text += delta['text']
                handler.markdown(text.replace("$", "USD ").replace("%", " percent"))

        elif "metadata" in chunk:
            st.session_state['input_token'] = chunk['metadata']['usage']["inputTokens"]
            st.session_state['output_token'] = chunk['metadata']['usage']["outputTokens"]
            latency = chunk['metadata']['metrics']["latencyMs"]
            pricing = st.session_state['input_token'] * pricing_file[f"anthropic.{params['model']}"]["input"] + st.session_state['output_token'] * pricing_file[f"anthropic.{params['model']}"]["output"]
            st.session_state['cost'] += pricing
    return text

def bedrock_claude_(params, chat_history, system_message, prompt, model_id, image_path=None, handler=None):
    chat_history_copy = chat_history[:]
    content = []

    # Xử lý hình ảnh (nếu có)
    if image_path:
        if not isinstance(image_path, list):
            image_path = [image_path]
        for img in image_path:
            s3 = boto3.client('s3', region_name=REGION)
            match = re.match("s3://(.+?)/(.+)", img)
            image_name = os.path.basename(img)
            _, ext = os.path.splitext(image_name)
            if "jpg" in ext:
                ext = ".jpeg"
            bucket_name = match.group(1)
            key = match.group(2)
            obj = s3.get_object(Bucket=bucket_name, Key=key)
            bytes_image = obj['Body'].read()
            content.extend([{"text": image_name}, {
                "image": {
                    "format": f"{ext.lower().replace('.', '')}",
                    "source": {"bytes": bytes_image}
                }
            }])

    content.append({
        "text": prompt
    })
    chat_history_copy.append({"role": "user", "content": content})
    system_message = [{"text": system_message}]

    # Kiểm tra mô hình và sử dụng API phù hợp
    if model_id == 'anthropic.claude-3-haiku-20240307-v1:0':
        # Sử dụng converse_stream cho claude-3-haiku
        response = bedrock_runtime.converse_stream(
            messages=chat_history_copy,
            modelId=model_id,
            inferenceConfig={"maxTokens": OUTPUT_TOKEN, "temperature": 0.5},
            system=system_message
        )
        answer = bedrock_streemer(params, response, handler)
    else:
        # Giữ nguyên logic cho các mô hình khác
        response = bedrock_runtime.converse_stream(
            messages=chat_history_copy,
            modelId=model_id,
            inferenceConfig={"maxTokens": OUTPUT_TOKEN, "temperature": 0.5},
            system=system_message
        )
        answer = bedrock_streemer(params, response, handler)

    return answer

def _invoke_bedrock_with_retries(params, current_chat, chat_template, question, model_id, image_path, handler):
    max_retries = 10
    backoff_base = 2
    max_backoff = 3  # Maximum backoff time in seconds
    retries = 0

    while True:
        try:
            response = bedrock_claude_(params, current_chat, chat_template, question, model_id, image_path, handler)
            return response
        except ClientError as e:
            if e.response['Error']['Code'] == 'ThrottlingException':
                if retries < max_retries:
                    # Throttling, exponential backoff
                    sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                    time.sleep(sleep_time)
                    retries += 1
                else:
                    raise e
            elif e.response['Error']['Code'] == 'ModelStreamErrorException':
                if retries < max_retries:
                    # Throttling, exponential backoff
                    sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                    time.sleep(sleep_time)
                    retries += 1
                else:
                    raise e
            elif e.response['Error']['Code'] == 'EventStreamError':
                if retries < max_retries:
                    # Throttling, exponential backoff
                    sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                    time.sleep(sleep_time)
                    retries += 1
                else:
                    raise e
            else:
                # Some other API error, rethrow
                raise

def get_session_ids_by_user(user_id):
    """
    Get Session Ids and corresponding top message for a user to populate the chat history drop down on the front end
    """
    message_list = {}
    cursor.execute("""
        SELECT "SessionID", "Messages" FROM "RangDongChatbot" WHERE "UserID"=%s
    """, (user_id,))
    results = cursor.fetchall()
    for session_id, messages in results:
        if messages:
            message_list[session_id] = messages[0]['user']
    return message_list

def list_csv_xlsx_in_s3_folder(bucket_name, folder_path):
    """
    List all CSV and XLSX files in a specified S3 folder.
    """
    s3 = boto3.client('s3')
    csv_xlsx_files = []

    try:
        # Ensure the folder path ends with a '/'
        if not folder_path.endswith('/'):
            folder_path += '/'

        # List objects in the specified folder
        paginator = s3.get_paginator('list_objects_v2')
        page_iterator = paginator.paginate(Bucket=bucket_name, Prefix=folder_path)

        for page in page_iterator:
            if 'Contents' in page:
                for obj in page['Contents']:
                    # Get the file name
                    file_name = obj['Key']

                    # Check if the file is a CSV or XLSX
                    if file_name.lower().endswith(INPUT_EXT):
                        csv_xlsx_files.append(os.path.basename(file_name))

        return csv_xlsx_files

    except ClientError as e:
        print(f"An error occurred: {e}")
        return []

def query_llm(params, handler):
    """
    Function takes a user query and an uploaded document. Caches documents in S3
    passing a document is optional
    """

    if not isinstance(params['upload_doc'], list):
        raise TypeError("documents must be in a list format")
     # Kiểm tra mô hình và thiết lập model_id
    if "claude" in params['model']:
        if params['model'] == 'claude-3-haiku':
            model_id = 'anthropic.claude-3-haiku-20240307-v1:0'
            claude3 = True
        else:
            model_id = 'anthropic.' + params['model']
            claude3 = False
    else:
        model_id = params['model'] + "-instruct-v1:0"
        claude3 = False
    # Check if Claude3 model is used and handle images with the CLAUDE3 Model
    # claude3 = False
    if "claude" in params['model']:
        model = 'anthropic.' + params['model']
    else:
        model = params['model'] + "-instruct-v1:0"

    if "sonnet" in model or "haiku" in model:
        model += "-20240620-v1:0" if "claude-3-5" in model else "-20240307-v1:0" if "haiku" in model else "-20240229-v1:0"
        claude3 = True

    ## prompt template for when a user uploads a doc
    doc_path = []
    image_path = []
    full_doc_path = []
    doc = ""
    if params['tools']:
        messages, tool, results, image_holder, doc_list, stop_reason = function_calling_utils.function_caller_claude_(params, handler)
        if stop_reason != "tool_use":
            return messages
        elif stop_reason == "tool_use":
            prompt = f"""You are a conversational AI Assistant. 
I will provide you with a question on a dataset, a python code that implements the solution to the question and the result of that code solution.
Here is the question:
<question>
{params['question']}
</question>

Here is the python code:
<python>
{tool['input']['code']}
</python>

Here the result of the code:
<result>
{results}
</result>

After reading the user question, respond with a detailed analytical answer based entirely on the result from the code. Do NOT make up answers. 
When providing your response:
- Do not include any preamble, go straight to the answer.
- It should not be obvious you are referencing the result."""

            system_message = "You always provide your response in a well presented format using markdown. Make use of tables, list etc. where necessary in your response, so information is well presented and easily read."
            answer = _invoke_bedrock_with_retries(params, [], system_message, prompt, model, image_holder, handler)

            chat_history = {"user": results["text"] if "text" in results else "",
                            "assistant": answer,
                            "image": image_holder,
                            "document": [],
                            "modelID": model,
                            "code": tool['input']['code'],
                            "time": str(time.time()),
                            "input_token": round(st.session_state['input_token']),
                            "output_token": round(st.session_state['output_token']),
                            "tool_result_id": tool['toolUseId'],
                            "tool_name": '',
                            "tool_params": ''}

            # Store conversation memory in PostgreSQL
            put_db(params, chat_history)
            return answer
    else:
        current_chat, chat_hist = get_chat_history_db(params, CHAT_HISTORY_LENGTH, claude3)
        if params['upload_doc'] or params['s3_objects']:
            if params['upload_doc']:
                doc = 'I have provided documents and/or images.\n'
                for ids, docs in enumerate(params['upload_doc']):
                    file_name = docs.name
                    _, extensions = os.path.splitext(file_name)
                    docs = put_obj_in_s3_bucket_(docs)
                    full_doc_path.append(docs)
                    if extensions.lower() in [".jpg", ".jpeg", ".png", ".gif", ".webp"] and claude3:
                        image_path.append(docs)
                        continue

            if params['s3_objects']:
                doc = 'I have provided documents and/or images.\n'
                for ids, docs in enumerate(params['s3_objects']):
                    file_name = docs
                    _, extensions = os.path.splitext(file_name)
                    docs = put_obj_in_s3_bucket_(f"s3://{INPUT_BUCKET}/{INPUT_S3_PATH}/{docs}")
                    full_doc_path.append(docs)
                    if extensions.lower() in [".jpg", ".jpeg", ".png", ".gif", ".webp"] and claude3:
                        image_path.append(docs)
                        continue

            doc_path = [item for item in full_doc_path if item not in image_path]
            errors, result_string = process_files(doc_path)
            if errors:
                st.error(errors)
            doc += result_string
            with open("prompt/doc_chat.txt", "r", encoding="utf-8") as f:
                chat_template = f.read()
        else:
            # Chat template for open ended query
            with open("prompt/chat.txt", "r", encoding="utf-8") as f:
                chat_template = f.read()
        response = _invoke_bedrock_with_retries(params, current_chat, chat_template, doc + params['question'], model, image_path, handler)
        # Log the following items to PostgreSQL
        chat_history = {"user": params['question'],
                        "assistant": response,
                        "image": image_path,
                        "document": doc_path,
                        "modelID": model,
                        "time": str(time.time()),
                        "input_token": round(st.session_state['input_token']),
                        "output_token": round(st.session_state['output_token'])}
        # Store conversation memory in PostgreSQL
        put_db(params, chat_history)
        return response

def get_chat_historie_for_streamlit(params):
    """
    This function retrieves chat history stored in PostgreSQL partitioned by a userID and SessionID
    """
    cursor.execute("""
        SELECT "Messages" FROM "RangDongChatbot" WHERE "UserID"=%s AND "SessionID"=%s
    """, (st.session_state['userid'], params["session_id"]))
    result = cursor.fetchone()
    if result:
        chat_histories = result[0]
    else:
        chat_histories = []

    # Constructing the desired list of dictionaries
    formatted_data = []
    if chat_histories:
        for entry in chat_histories:
            image_files = [os.path.basename(x) for x in entry.get('image', [])]
            doc_files = [os.path.basename(x) for x in entry.get('document', [])]
            code_script = entry.get('code', "")
            assistant_attachment = '\n\n'.join(image_files + doc_files)
            if "tool_result_id" in entry and not entry["tool_result_id"]:
                formatted_data.append({
                    "role": "user",
                    "content": entry["user"],
                })
            elif "tool_result_id" not in entry:
                formatted_data.append({
                    "role": "user",
                    "content": entry["user"],
                })
            if "tool_use_id" in entry and not entry["tool_use_id"]:
                formatted_data.append({
                    "role": "assistant",
                    "content": entry["assistant"],
                    "attachment": assistant_attachment,
                    "code": code_script,
                })
            elif "tool_use_id" not in entry:
                formatted_data.append({
                    "role": "assistant",
                    "content": entry["assistant"],
                    "attachment": assistant_attachment,
                    "code": code_script,
                    "code-result": entry["user"],
                    "image_output": entry.get('image', []) if "tool_result_id" in entry else []
                })
    else:
        chat_histories = []
    return formatted_data, chat_histories


def get_key_from_value(dictionary, value):
    return next((key for key, val in dictionary.items() if val == value), None)

def chat_bedrock_(params):
    st.title('Chatbot AI Assistant 🙂')
    params['chat_histories'] = []
    if params["session_id"].strip():
        st.session_state.messages, params['chat_histories'] = get_chat_historie_for_streamlit(params)
    for message in st.session_state.messages:

        with st.chat_message(message["role"]):
            if "```" in message["content"]:
                st.markdown(message["content"], unsafe_allow_html=True)
            else:
                st.markdown(message["content"].replace("$", r"\$"), unsafe_allow_html=True)
            if message["role"] == "assistant":
                if message.get("image_output"):
                    for item in message["image_output"]:
                        bucket_name, key = item.replace('s3://', '').split('/', 1)
                        image_bytes = get_object_with_retry(bucket_name, key)
                        image = Image.open(io.BytesIO(image_bytes['Body'].read()))
                        st.image(image)
                if message.get("attachment"):
                    with st.expander(label="**attachments**"):
                        st.markdown(message["attachment"])
                if message.get('code'):
                    with st.expander(label="**code snippet**"):
                        st.markdown(f'```python\n{message["code"]}', unsafe_allow_html=True)
                    with st.expander(label="**code result**"):
                        st.markdown(f'```python\n{message["code-result"]}', unsafe_allow_html=True)

    if prompt := st.chat_input("What's up?"):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt.replace("$", r"\$"), unsafe_allow_html=True)
        with st.chat_message("assistant"):
            message_placeholder = st.empty()
            params["question"] = prompt
            answer = query_llm(params, message_placeholder)
            message_placeholder.markdown(answer.replace("$", r"\$"), unsafe_allow_html=True)
            st.session_state.messages.append({"role": "assistant", "content": answer})
        st.rerun()

def app_sidebar():
    with st.sidebar:
        st.metric(label="Bedrock Session Cost", value=f"${round(st.session_state['cost'], 2)}")
        st.write("-----")
        button = st.button("New Chat", type="primary")
        models = ['claude-3-5-sonnet', 'claude-3-sonnet', 'claude-3-haiku', 'claude-instant-v1', 'claude-v2:1', 'claude-v2']
        model = st.selectbox('**Model**', models)
        params = {"model": model}
        user_sess_id = get_session_ids_by_user(st.session_state['userid'])
        float_keys = {float(key): value for key, value in user_sess_id.items()}
        sorted_messages = sorted(float_keys.items(), reverse=True)
        sorted_messages.insert(0, (float(st.session_state['user_sess']), "New Chat"))
        if button:
            st.session_state['user_sess'] = str(time.time())
            sorted_messages.insert(0, (float(st.session_state['user_sess']), "New Chat"))
        st.session_state['chat_session_list'] = dict(sorted_messages)
        chat_items = st.selectbox("**Chat Sessions**", st.session_state['chat_session_list'].values(), key="chat_sessions")
        session_id = get_key_from_value(st.session_state['chat_session_list'], chat_items)
        tools = st.multiselect("**Tools**", ["Advanced Data Analytics"], key="function_collen", default=None)
        bucket_items = list_csv_xlsx_in_s3_folder(INPUT_BUCKET, INPUT_S3_PATH)
        bucket_objects = st.multiselect("**Files**", bucket_items, key="objector", default=None)
        file = st.file_uploader('Upload a document', accept_multiple_files=True, help="pdf,csv,txt,png,jpg,xlsx,json,py doc format supported")
        if file and LOAD_DOC_IN_ALL_CHAT_CONVO:
            st.warning('You have set **load-doc-in-chat-history** to true. For better performance, remove uploaded file(s) (by clicking **X**) **AFTER** first query on uploaded files. See the README for more info', icon="⚠️")
        params = {"model": model, "session_id": str(session_id), "chat_item": chat_items, "upload_doc": file, "tools": tools, 's3_objects': bucket_objects}
        st.session_state['count'] = 1
        return params

def main():
    params = app_sidebar()
    chat_bedrock_(params)

if __name__ == '__main__':
    main()
